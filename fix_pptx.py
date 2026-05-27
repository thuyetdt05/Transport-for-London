# -*- coding: utf-8 -*-
"""
Script sửa Nhom14_TfL_ETL_Pipeline.pptx dựa trên kết quả kiểm tra đồng bộ.

3 điểm sửa:
  [1] SLIDE 8  — Sửa Features KMeans (TextBox 84)
  [2] SLIDE 7  — Sửa Spatial Join library (TextBox 20 / ID 21)
  [3] SLIDE 10 — Cập nhật Top 5 ga đúng số liệu từ CSV
"""

import sys
import copy
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

SRC = Path("Nhom14_TfL_ETL_Pipeline.pptx")
DST = Path("Nhom14_TfL_ETL_Pipeline.pptx")  # overwrite in-place


def get_shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_text_keep_format(tf, new_text):
    """Ghi đè nội dung text frame, giữ nguyên font/size/color của paragraph đầu."""
    # Lấy format từ run đầu tiên nếu có
    para0 = tf.paragraphs[0]
    ref_run = para0.runs[0] if para0.runs else None

    # Xóa hết paragraph, chỉ giữ 1
    while len(tf.paragraphs) > 1:
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)

    # Ghi nội dung mới vào paragraph đầu
    p = tf.paragraphs[0]
    # Xóa hết run
    for r in p.runs:
        r._r.getparent().remove(r._r)

    run = p.add_run()
    run.text = new_text

    # Copy format từ ref_run nếu có
    if ref_run:
        if ref_run.font.bold is not None:
            run.font.bold = ref_run.font.bold
        if ref_run.font.size:
            run.font.size = ref_run.font.size
        if ref_run.font.color and ref_run.font.color.type:
            try:
                run.font.color.rgb = ref_run.font.color.rgb
            except Exception:
                pass
        if ref_run.font.name:
            run.font.name = ref_run.font.name


def fix_slide7(slide):
    """Sửa 'GeoPandas sjoin_nearest()' → 'Pandas + NumPy Nearest Neighbor'."""
    shape = get_shape_by_name(slide, "TextBox 20")  # ID 21
    if shape and shape.has_text_frame:
        tf = shape.text_frame
        old_text = tf.paragraphs[0].text
        new_text = "▸  Pandas + NumPy Euclidean Nearest Neighbor – không cần GeoPandas"
        set_text_keep_format(tf, new_text)
        print(f"  [Slide 7] TextBox 20: '{old_text}' → '{new_text}'")

    # Sửa thêm dòng Tolerance (không còn phù hợp với pandas approach)
    shape2 = get_shape_by_name(slide, "TextBox 23")  # "Tolerance: 500m..."
    if shape2 and shape2.has_text_frame:
        tf2 = shape2.text_frame
        old2 = tf2.paragraphs[0].text
        new2 = "▸  Ưu tiên StopType (MET > RLY > DLR) + khoảng cách tọa độ gần nhất"
        set_text_keep_format(tf2, new2)
        print(f"  [Slide 7] TextBox 23: '{old2}' → '{new2}'")


def fix_slide8(slide):
    """Sửa Features KMeans từ avg_passengers,covid... → passengers_2021,num_lines,lat,lon."""
    shape = get_shape_by_name(slide, "TextBox 84")  # ID 85
    if shape and shape.has_text_frame:
        tf = shape.text_frame
        old_text = tf.paragraphs[0].text
        new_text = (
            "📐  Features: passengers_2021, num_lines, lat, lon  |  "
            "StandardScaler → KMeans (sklearn)  |  Elbow Method → K=6"
        )
        set_text_keep_format(tf, new_text)
        print(f"  [Slide 8] TextBox 84: Features đã được cập nhật.")
        print(f"    Cũ: {old_text}")
        print(f"    Mới: {new_text}")


# Số liệu TOP 5 đúng từ CSV (passengers_2021 thực tế)
TOP5_CORRECT = [
    ("Stratford",                 "63.4M"),
    ("Liverpool Street",          "43.1M"),
    ("King's Cross St. Pancras",  "36.7M"),
    ("Victoria",                  "33.5M"),
    ("Oxford Circus",             "32.9M"),
]

# Shape names tương ứng với hạng 1..5 trong Slide 10
# Mỗi hạng có 3 shape: số thứ tự (không đổi), tên ga, số HK
TOP5_SHAPE_MAP = {
    1: {"name": "TextBox 27", "num": "TextBox 29"},
    2: {"name": "TextBox 31", "num": "TextBox 33"},
    3: {"name": "TextBox 35", "num": "TextBox 37"},
    4: {"name": "TextBox 39", "num": "TextBox 41"},
    5: {"name": "TextBox 43", "num": "TextBox 45"},
}


def fix_slide10(slide):
    """Cập nhật Top 5 ga đúng theo số liệu CSV."""
    for rank, (station_name, station_num) in enumerate(TOP5_CORRECT, start=1):
        map_entry = TOP5_SHAPE_MAP[rank]

        # Sửa tên ga
        shape_name = get_shape_by_name(slide, map_entry["name"])
        if shape_name and shape_name.has_text_frame:
            old = shape_name.text_frame.paragraphs[0].text
            set_text_keep_format(shape_name.text_frame, station_name)
            print(f"  [Slide 10] #{rank} Tên: '{old}' → '{station_name}'")

        # Sửa số HK
        shape_num = get_shape_by_name(slide, map_entry["num"])
        if shape_num and shape_num.has_text_frame:
            old = shape_num.text_frame.paragraphs[0].text
            set_text_keep_format(shape_num.text_frame, station_num)
            print(f"  [Slide 10] #{rank} Số HK: '{old}' → '{station_num}'")


# ─── MAIN ────────────────────────────────────────────────────────────────────

def main():
    print("=" * 65)
    print("SỬA BÀI TRÌNH CHIẾU: Nhom14_TfL_ETL_Pipeline.pptx")
    print("=" * 65)

    prs = Presentation(str(SRC))

    print("\n[1] SLIDE 7 — Sửa Spatial Join (GeoPandas → Pandas+NumPy)")
    fix_slide7(prs.slides[6])   # 0-indexed

    print("\n[2] SLIDE 8 — Sửa Features KMeans")
    fix_slide8(prs.slides[7])

    print("\n[3] SLIDE 10 — Cập nhật Top 5 ga đúng số liệu")
    fix_slide10(prs.slides[9])

    prs.save(str(DST))
    size_kb = DST.stat().st_size / 1024
    print(f"\n{'=' * 65}")
    print(f"✅ Đã lưu: {DST}  ({size_kb:.1f} KB)")
    print("=" * 65)


if __name__ == "__main__":
    main()
