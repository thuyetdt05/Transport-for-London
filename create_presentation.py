"""
Tạo file PowerPoint chuyên nghiệp cho Luận văn Nhóm 14 - DHKHDL20A
Dự án: Xây dựng Pipeline ETL End-to-End & Phân tích Dữ liệu Không gian - Thời gian
Hệ thống Nhà ga TfL London
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from pptx.util import Emu
import copy
from lxml import etree
import os

# ── Đường dẫn ──────────────────────────────────────────────────────────────
BASE_DIR   = r"c:\Users\My PC\Downloads\DE"
OUTPUT_DIR = os.path.join(BASE_DIR, "outputs")
IMG_CLUSTER  = os.path.join(OUTPUT_DIR, "chart_cluster_comparison.png")
IMG_TRENDS   = os.path.join(OUTPUT_DIR, "chart_passenger_trends.png")
OUT_PPTX     = os.path.join(BASE_DIR, "Nhom14_TfL_ETL_Pipeline.pptx")

# ── Bảng màu ───────────────────────────────────────────────────────────────
TEAL       = RGBColor(0x00, 0xC4, 0xB4)   # #00C4B4
DARK_NAVY  = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A
MID_NAVY   = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B  (slide bgs)
SLATE      = RGBColor(0x33, 0x4E, 0x68)   # #334E68
LIGHT_GREY = RGBColor(0xF1, 0xF5, 0xF9)   # #F1F5F9
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_YLW = RGBColor(0xFF, 0xC1, 0x07)   # #FFC107  amber accent
ACCENT_BLU = RGBColor(0x38, 0xBD, 0xF8)   # #38BDF8  sky-blue

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════
# Helper utilities
# ═══════════════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Add a completely blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def fill_slide_bg(slide, color: RGBColor):
    """Set a solid background colour for a single slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name="Calibri", font_size=18,
                bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT,
                wrap=True, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_name="Calibri", font_size=16,
             bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             space_before_pt=4, level=0):
    """Append a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before_pt)
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return p


def add_image(slide, img_path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        slide.shapes.add_picture(img_path, left, top, width)


def teal_accent_bar(slide, top=Inches(0.82), height=Inches(0.055)):
    """Full-width teal accent line just below the header zone."""
    add_rect(slide, 0, top, SLIDE_W, height, fill_color=TEAL)


def header_band(slide, title_text, subtitle_text="", dark_bg=True):
    """Dark top band with title + optional subtitle."""
    band_h = Inches(1.0)
    band_color = DARK_NAVY if dark_bg else MID_NAVY
    add_rect(slide, 0, 0, SLIDE_W, band_h, fill_color=band_color)
    teal_accent_bar(slide, top=band_h - Inches(0.07), height=Inches(0.07))

    # Title
    add_textbox(slide, title_text,
                left=Inches(0.45), top=Inches(0.12),
                width=Inches(10.5), height=Inches(0.65),
                font_name="Calibri", font_size=26, bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    left=Inches(0.45), top=Inches(0.68),
                    width=Inches(10), height=Inches(0.3),
                    font_name="Calibri", font_size=13, bold=False,
                    color=RGBColor(0xA0, 0xAE, 0xC0), align=PP_ALIGN.LEFT)


def slide_number_label(slide, num, total=13):
    add_textbox(slide, f"{num} / {total}",
                left=Inches(12.3), top=Inches(7.1),
                width=Inches(0.9), height=Inches(0.3),
                font_name="Calibri", font_size=10,
                color=RGBColor(0x64, 0x74, 0x8B),
                align=PP_ALIGN.RIGHT)


def footer_bar(slide, text="Nhóm 14 – DHKHDL20A  |  GVHD: TS. Lê Trọng Ngọc"):
    add_rect(slide, 0, Inches(7.22), SLIDE_W, Inches(0.28), fill_color=DARK_NAVY)
    add_textbox(slide, text,
                left=Inches(0.3), top=Inches(7.23),
                width=Inches(10), height=Inches(0.25),
                font_name="Calibri", font_size=9,
                color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)


def info_card(slide, left, top, width, height,
              icon_char, label, value,
              bg_color=MID_NAVY, accent=TEAL):
    """Small KPI card with icon, label, value."""
    add_rect(slide, left, top, width, height,
             fill_color=bg_color, line_color=accent, line_width_pt=1.5)
    # icon
    add_textbox(slide, icon_char,
                left=left + Inches(0.12), top=top + Inches(0.1),
                width=Inches(0.5), height=Inches(0.4),
                font_name="Segoe UI Emoji", font_size=22, color=TEAL)
    # label
    add_textbox(slide, label,
                left=left + Inches(0.12), top=top + Inches(0.48),
                width=width - Inches(0.2), height=Inches(0.28),
                font_name="Calibri", font_size=11,
                color=RGBColor(0x94, 0xA3, 0xB8))
    # value
    add_textbox(slide, value,
                left=left + Inches(0.12), top=top + Inches(0.75),
                width=width - Inches(0.2), height=Inches(0.38),
                font_name="Calibri", font_size=18, bold=True, color=WHITE)


def bullet_box(slide, items, left, top, width, height,
               font_size=15, header=None, header_size=17,
               bg=MID_NAVY, border=TEAL):
    """Rounded-corner look: box + bulleted text."""
    add_rect(slide, left, top, width, height,
             fill_color=bg, line_color=border, line_width_pt=1.2)

    y_offset = Inches(0.08)
    if header:
        add_textbox(slide, header,
                    left=left + Inches(0.18), top=top + y_offset,
                    width=width - Inches(0.3), height=Inches(0.38),
                    font_name="Calibri", font_size=header_size,
                    bold=True, color=TEAL)
        y_offset += Inches(0.4)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.18), top + y_offset,
        width - Inches(0.3), height - y_offset - Inches(0.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = item
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.color.rgb = LIGHT_GREY


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def slide01_cover(prs):
    """Slide 1 – Trang bìa"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_NAVY)

    # Left dark panel
    add_rect(slide, 0, 0, Inches(7.2), SLIDE_H, fill_color=MID_NAVY)

    # Teal top accent stripe
    add_rect(slide, 0, 0, Inches(7.2), Inches(0.09), fill_color=TEAL)

    # TfL-style red circle logo area (decorative)
    circle = slide.shapes.add_shape(9, Inches(0.38), Inches(0.55), Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = RGBColor(0xDC, 0x14, 0x1C)
    circle.line.fill.background()

    add_textbox(slide, "TfL",
                left=Inches(0.44), top=Inches(0.6),
                width=Inches(0.6), height=Inches(0.5),
                font_name="Calibri", font_size=14, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    # Project tag
    add_textbox(slide, "ĐỒ ÁN MÔN HỌC  |  DATA ENGINEERING",
                left=Inches(0.38), top=Inches(1.45),
                width=Inches(6.5), height=Inches(0.35),
                font_name="Calibri", font_size=11, bold=False,
                color=TEAL, align=PP_ALIGN.LEFT)

    # Main title line 1
    add_textbox(slide, "XÂY DỰNG PIPELINE ETL",
                left=Inches(0.38), top=Inches(1.9),
                width=Inches(6.5), height=Inches(0.9),
                font_name="Calibri", font_size=36, bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)

    # Main title line 2
    add_textbox(slide, "END-TO-END &",
                left=Inches(0.38), top=Inches(2.75),
                width=Inches(6.5), height=Inches(0.75),
                font_name="Calibri", font_size=36, bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, "PHÂN TÍCH DỮ LIỆU KHÔNG GIAN",
                left=Inches(0.38), top=Inches(3.45),
                width=Inches(6.5), height=Inches(0.75),
                font_name="Calibri", font_size=28, bold=True,
                color=TEAL, align=PP_ALIGN.LEFT)

    add_textbox(slide, "THỜI GIAN",
                left=Inches(0.38), top=Inches(4.12),
                width=Inches(6.5), height=Inches(0.6),
                font_name="Calibri", font_size=28, bold=True,
                color=TEAL, align=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, "Hệ thống Nhà ga Transport for London (TfL)",
                left=Inches(0.38), top=Inches(4.8),
                width=Inches(6.5), height=Inches(0.38),
                font_name="Calibri", font_size=15, italic=True,
                color=RGBColor(0xA0, 0xAE, 0xC0), align=PP_ALIGN.LEFT)

    # Divider line
    add_rect(slide, Inches(0.38), Inches(5.28), Inches(2.5), Inches(0.04), fill_color=TEAL)

    # Group info
    info = [
        "👥  Nhóm 14 – DHKHDL20A",
        "👨‍🏫  GVHD: TS. Lê Trọng Ngọc",
        "📅  Năm học: 2024 – 2025",
    ]
    y = Inches(5.42)
    for line in info:
        add_textbox(slide, line,
                    left=Inches(0.38), top=y,
                    width=Inches(6.5), height=Inches(0.34),
                    font_name="Calibri", font_size=13,
                    color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.LEFT)
        y += Inches(0.38)

    # Right panel - decorative map graphic simulation
    add_rect(slide, Inches(7.2), 0, Inches(6.13), SLIDE_H, fill_color=DARK_NAVY)

    # Grid lines decoration (simulate map)
    for i in range(8):
        add_rect(slide, Inches(7.2), Inches(i * 0.94), Inches(6.13), Inches(0.01),
                 fill_color=RGBColor(0x1E, 0x29, 0x3B))
    for i in range(7):
        add_rect(slide, Inches(7.2 + i * 0.88), 0, Inches(0.01), SLIDE_H,
                 fill_color=RGBColor(0x1E, 0x29, 0x3B))

    # Central large decorative circle (London Underground roundel style)
    big_circle = slide.shapes.add_shape(9,
        Inches(8.5), Inches(1.8), Inches(3.2), Inches(3.2))
    big_circle.fill.solid()
    big_circle.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    big_circle.line.color.rgb = TEAL
    big_circle.line.width = Pt(3)

    # Inner circle
    inner_c = slide.shapes.add_shape(9,
        Inches(9.0), Inches(2.3), Inches(2.2), Inches(2.2))
    inner_c.fill.solid()
    inner_c.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    inner_c.line.color.rgb = RGBColor(0x38, 0xBD, 0xF8)
    inner_c.line.width = Pt(1.5)

    # Horizontal bar through circle (roundel)
    add_rect(slide, Inches(8.3), Inches(3.2), Inches(3.6), Inches(0.45),
             fill_color=RGBColor(0xDC, 0x14, 0x1C))
    add_textbox(slide, "LONDON UNDERGROUND",
                left=Inches(8.3), top=Inches(3.24),
                width=Inches(3.6), height=Inches(0.38),
                font_name="Calibri", font_size=9, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "🚇",
                left=Inches(9.6), top=Inches(2.5),
                width=Inches(0.9), height=Inches(0.8),
                font_name="Segoe UI Emoji", font_size=36, color=TEAL)

    # Stats callouts on right panel
    stats = [
        ("298", "Nhà ga"),
        ("5 năm", "2017 – 2021"),
        ("3", "Nguồn dữ liệu"),
        ("K=6", "Cụm KMeans"),
    ]
    for idx, (val, lbl) in enumerate(stats):
        col = idx % 2
        row = idx // 2
        cx = Inches(7.6 + col * 2.9)
        cy = Inches(5.3 + row * 0.9)
        add_rect(slide, cx, cy, Inches(2.5), Inches(0.75),
                 fill_color=RGBColor(0x1E, 0x29, 0x3B),
                 line_color=TEAL, line_width_pt=1)
        add_textbox(slide, val,
                    left=cx + Inches(0.08), top=cy + Inches(0.04),
                    width=Inches(1.2), height=Inches(0.38),
                    font_name="Calibri", font_size=20, bold=True, color=TEAL)
        add_textbox(slide, lbl,
                    left=cx + Inches(0.08), top=cy + Inches(0.4),
                    width=Inches(2.3), height=Inches(0.28),
                    font_name="Calibri", font_size=10,
                    color=RGBColor(0x94, 0xA3, 0xB8))


# ───────────────────────────────────────────────────────────────────────────

def slide02_intro(prs):
    """Slide 2 – Lời mở đầu & Tóm tắt"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, LIGHT_GREY)
    header_band(slide, "LỜI MỞ ĐẦU & TÓM TẮT DỰ ÁN",
                "Bối cảnh, mục tiêu và kết quả tổng quan")
    footer_bar(slide)
    slide_number_label(slide, 2)

    # Left column – context + objectives
    add_rect(slide, Inches(0.3), Inches(1.15), Inches(5.9), Inches(5.85),
             fill_color=WHITE, line_color=TEAL, line_width_pt=1.5)

    add_textbox(slide, "📌  BỐI CẢNH & MỤC TIÊU",
                left=Inches(0.48), top=Inches(1.25),
                width=Inches(5.5), height=Inches(0.4),
                font_name="Calibri", font_size=15, bold=True, color=DARK_NAVY)

    context_items = [
        "• TfL London vận hành mạng lưới giao thông lớn nhất châu Âu với hàng triệu hành khách mỗi ngày.",
        "• Bùng phát COVID-19 (2020) gây biến động lớn về lưu lượng hành khách tại tất cả các ga.",
        "• Nhu cầu phân tích không gian – thời gian để đưa ra quyết định vận hành chính xác.",
        "",
        "🎯  Mục tiêu:",
        "  ➤  Xây dựng Pipeline ETL tự động hóa hoàn chỉnh",
        "  ➤  Tích hợp 3 nguồn dữ liệu dị biệt thành CSDL thống nhất",
        "  ➤  Phân cụm (KMeans) hành vi hành khách theo không gian",
        "  ➤  Phân tích tác động COVID-19 và xu hướng phục hồi",
        "  ➤  Trực quan hóa tương tác (Folium Map, biểu đồ)",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.48), Inches(1.75),
                                      Inches(5.55), Inches(5.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in context_items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.name = "Calibri"
        run.font.size = Pt(13)
        if item.startswith("🎯") or item.startswith("  ➤"):
            run.font.color.rgb = DARK_NAVY
            run.font.bold = item.startswith("🎯")
        else:
            run.font.color.rgb = SLATE

    # Right column – KPIs + scope
    cards = [
        ("🚉", "Nhà ga", "298 ga"),
        ("📅", "Giai đoạn", "2017 – 2021"),
        ("🗄️", "Nguồn dữ liệu", "3 nguồn"),
        ("🤖", "Thuật toán", "KMeans K=6"),
        ("🗺️", "Sản phẩm", "Folium Map"),
        ("📊", "Kết quả", "-63.21% COVID"),
    ]
    card_w = Inches(3.3)
    card_h = Inches(1.2)
    start_x = Inches(6.65)
    start_y = Inches(1.15)
    gap_x = Inches(0.22)
    gap_y = Inches(0.18)

    for i, (icon, lbl, val) in enumerate(cards):
        col = i % 2
        row = i // 2
        cx = start_x + col * (card_w + gap_x)
        cy = start_y + row * (card_h + gap_y)
        add_rect(slide, cx, cy, card_w, card_h,
                 fill_color=DARK_NAVY, line_color=TEAL, line_width_pt=1.5)
        add_textbox(slide, icon,
                    left=cx + Inches(0.15), top=cy + Inches(0.12),
                    width=Inches(0.5), height=Inches(0.45),
                    font_name="Segoe UI Emoji", font_size=22, color=TEAL)
        add_textbox(slide, lbl,
                    left=cx + Inches(0.65), top=cy + Inches(0.12),
                    width=Inches(2.4), height=Inches(0.3),
                    font_name="Calibri", font_size=11,
                    color=RGBColor(0x94, 0xA3, 0xB8))
        add_textbox(slide, val,
                    left=cx + Inches(0.65), top=cy + Inches(0.45),
                    width=Inches(2.4), height=Inches(0.55),
                    font_name="Calibri", font_size=18, bold=True, color=WHITE)


# ───────────────────────────────────────────────────────────────────────────

def slide03_reasons(prs):
    """Slide 3 – Lý do chọn đề tài"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_NAVY)
    header_band(slide, "LÝ DO CHỌN ĐỀ TÀI", "Tính cấp thiết và giá trị thực tiễn của dự án")
    footer_bar(slide)
    slide_number_label(slide, 3)

    reasons = [
        ("🏙️", "Mạng lưới TfL cực kỳ phức tạp",
         "270+ tuyến, 11 loại dịch vụ, dữ liệu hành khách từ nhiều hệ thống khác nhau. Cần giải pháp ETL tích hợp toàn diện."),
        ("🦠", "Tác động sâu sắc của COVID-19",
         "Đại dịch 2020 gây sụt giảm hành khách -63.21% – cơ hội nghiên cứu hành vi di chuyển đô thị trên quy mô lớn."),
        ("🗺️", "Kết hợp Không gian & Thời gian",
         "Phân tích đồng thời tọa độ địa lý và xu hướng theo thời gian → bức tranh toàn diện về hoạt động đô thị."),
        ("📈", "Giá trị thực tiễn cao",
         "Hỗ trợ TfL tối ưu phân bổ nguồn lực, lập kế hoạch phục hồi và hoạch định chính sách giao thông dài hạn."),
    ]

    card_w = Inches(5.9)
    card_h = Inches(2.2)
    positions = [
        (Inches(0.3),  Inches(1.15)),
        (Inches(6.7),  Inches(1.15)),
        (Inches(0.3),  Inches(3.55)),
        (Inches(6.7),  Inches(3.55)),
    ]

    for (icon, title, body), (lx, ly) in zip(reasons, positions):
        add_rect(slide, lx, ly, card_w, card_h,
                 fill_color=MID_NAVY, line_color=TEAL, line_width_pt=1.8)
        # Icon circle
        circ = slide.shapes.add_shape(9, lx + Inches(0.18), ly + Inches(0.25),
                                       Inches(0.65), Inches(0.65))
        circ.fill.solid(); circ.fill.fore_color.rgb = DARK_NAVY
        circ.line.color.rgb = TEAL; circ.line.width = Pt(1.5)
        add_textbox(slide, icon,
                    left=lx + Inches(0.19), top=ly + Inches(0.28),
                    width=Inches(0.6), height=Inches(0.5),
                    font_name="Segoe UI Emoji", font_size=20, color=TEAL,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, title,
                    left=lx + Inches(0.95), top=ly + Inches(0.22),
                    width=card_w - Inches(1.1), height=Inches(0.42),
                    font_name="Calibri", font_size=15, bold=True, color=TEAL)
        add_textbox(slide, body,
                    left=lx + Inches(0.18), top=ly + Inches(0.78),
                    width=card_w - Inches(0.32), height=Inches(1.3),
                    font_name="Calibri", font_size=12.5,
                    color=RGBColor(0xCB, 0xD5, 0xE1))


# ───────────────────────────────────────────────────────────────────────────

def slide04_datasets(prs):
    """Slide 4 – Bộ dữ liệu sử dụng"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, LIGHT_GREY)
    header_band(slide, "BỘ DỮ LIỆU SỬ DỤNG",
                "3 nguồn dữ liệu độc lập được tích hợp trong Pipeline ETL")
    footer_bar(slide)
    slide_number_label(slide, 4)

    sources = [
        {
            "num": "01",
            "file": "stations.kml",
            "type": "Dữ liệu địa lý (KML)",
            "desc": "File KML từ Google Maps chứa tọa độ địa lý của các nhà ga TfL.",
            "fields": ["📍 Tên ga (name)", "🌐 Kinh độ / Vĩ độ", "📝 Mô tả HTML", "📋 ~298 điểm địa lý"],
            "color": RGBColor(0x06, 0x95, 0x88),
        },
        {
            "num": "02",
            "file": "TfL_stations.csv",
            "type": "Dữ liệu hành khách (CSV)",
            "desc": "Dữ liệu lưu lượng hành khách hàng năm do TfL công bố chính thức.",
            "fields": ["🚉 Tên ga", "📊 Lượt vào/ra 2017-2021", "🏙️ Borough", "📋 ~588 bản ghi"],
            "color": RGBColor(0x29, 0x78, 0xC8),
        },
        {
            "num": "03",
            "file": "Stops.csv",
            "type": "Dữ liệu tuyến dừng (CSV)",
            "desc": "Dữ liệu điểm dừng từ hệ thống lên kế hoạch chuyến đi của TfL.",
            "fields": ["🔑 ATCOCode / CommonName", "🌐 Tọa độ bổ sung", "🚌 Loại phương tiện", "📋 ~101MB"],
            "color": RGBColor(0x76, 0x44, 0xC2),
        },
    ]

    col_w = Inches(4.1)
    card_h = Inches(5.85)
    gap = Inches(0.22)
    start_x = Inches(0.32)
    card_y = Inches(1.15)

    for i, src in enumerate(sources):
        cx = start_x + i * (col_w + gap)
        add_rect(slide, cx, card_y, col_w, card_h,
                 fill_color=WHITE, line_color=src["color"], line_width_pt=2)

        # Number badge
        badge = slide.shapes.add_shape(9, cx + col_w - Inches(0.68), card_y + Inches(0.12),
                                        Inches(0.55), Inches(0.55))
        badge.fill.solid(); badge.fill.fore_color.rgb = src["color"]
        badge.line.fill.background()
        add_textbox(slide, src["num"],
                    left=cx + col_w - Inches(0.67), top=card_y + Inches(0.15),
                    width=Inches(0.52), height=Inches(0.42),
                    font_name="Calibri", font_size=14, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)

        # File name
        add_textbox(slide, src["file"],
                    left=cx + Inches(0.18), top=card_y + Inches(0.15),
                    width=col_w - Inches(0.85), height=Inches(0.45),
                    font_name="Calibri", font_size=17, bold=True,
                    color=DARK_NAVY)

        # Type badge
        add_rect(slide, cx + Inches(0.18), card_y + Inches(0.68),
                 col_w - Inches(0.36), Inches(0.3),
                 fill_color=src["color"])
        add_textbox(slide, src["type"],
                    left=cx + Inches(0.2), top=card_y + Inches(0.68),
                    width=col_w - Inches(0.36), height=Inches(0.3),
                    font_name="Calibri", font_size=10, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, src["desc"],
                    left=cx + Inches(0.18), top=card_y + Inches(1.1),
                    width=col_w - Inches(0.36), height=Inches(0.85),
                    font_name="Calibri", font_size=12, italic=True,
                    color=SLATE)

        # Divider
        add_rect(slide, cx + Inches(0.18), card_y + Inches(2.05),
                 col_w - Inches(0.36), Inches(0.03),
                 fill_color=RGBColor(0xE2, 0xE8, 0xF0))

        # Fields
        txBox = slide.shapes.add_textbox(cx + Inches(0.18), card_y + Inches(2.18),
                                          col_w - Inches(0.36), Inches(3.4))
        tf = txBox.text_frame; tf.word_wrap = True
        first = True
        for field in src["fields"]:
            if first:
                p = tf.paragraphs[0]; first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(7)
            run = p.add_run()
            run.text = field
            run.font.name = "Calibri"
            run.font.size = Pt(13)
            run.font.color.rgb = DARK_NAVY


# ───────────────────────────────────────────────────────────────────────────

def slide05_challenges(prs):
    """Slide 5 – Thách thức & Giải pháp"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_NAVY)
    header_band(slide, "THÁCH THỨC DỮ LIỆU THÔ & GIẢI PHÁP",
                "Vấn đề phát sinh trong quá trình tích hợp và cách xử lý")
    footer_bar(slide)
    slide_number_label(slide, 5)

    # Left: Challenges
    add_rect(slide, Inches(0.3), Inches(1.15), Inches(5.85), Inches(5.9),
             fill_color=RGBColor(0x1E, 0x1E, 0x2E), line_color=RGBColor(0xDC, 0x14, 0x1C), line_width_pt=2)

    add_textbox(slide, "⚠️  THÁCH THỨC",
                left=Inches(0.48), top=Inches(1.25),
                width=Inches(5.5), height=Inches(0.42),
                font_name="Calibri", font_size=16, bold=True,
                color=RGBColor(0xFF, 0x6B, 0x6B))

    challenges = [
        ("❌", "Tên ga không đồng nhất",
         "\"King's Cross\" vs \"Kings Cross\" vs \"King'S Cross St. Pancras\" – không thể join trực tiếp"),
        ("❌", "File Stops.csv kích thước lớn",
         "~101MB với >500,000 dòng – cần tối ưu bộ nhớ, đọc theo chunk"),
        ("❌", "Thiếu thông tin Borough",
         "Nhiều bản ghi KML không có Borough → cần suy luận từ tọa độ địa lý"),
        ("❌", "Dữ liệu hành khách trống/lỗi",
         "Năm đặc biệt 2020 có giá trị bất thường cần xử lý cẩn thận"),
        ("❌", "Định dạng dữ liệu không đồng nhất",
         "CSV, KML, database – mỗi định dạng cần parser riêng biệt"),
    ]

    y = Inches(1.75)
    for icon, title, desc in challenges:
        add_textbox(slide, f"{icon}  {title}",
                    left=Inches(0.48), top=y,
                    width=Inches(5.5), height=Inches(0.32),
                    font_name="Calibri", font_size=13, bold=True,
                    color=RGBColor(0xFF, 0xA5, 0xA5))
        add_textbox(slide, f"     {desc}",
                    left=Inches(0.48), top=y + Inches(0.3),
                    width=Inches(5.5), height=Inches(0.45),
                    font_name="Calibri", font_size=11,
                    color=RGBColor(0x94, 0xA3, 0xB8))
        y += Inches(0.93)

    # Arrow
    add_textbox(slide, "→",
                left=Inches(6.32), top=Inches(3.5),
                width=Inches(0.65), height=Inches(0.6),
                font_name="Calibri", font_size=32, bold=True,
                color=TEAL, align=PP_ALIGN.CENTER)

    # Right: Solutions
    add_rect(slide, Inches(7.1), Inches(1.15), Inches(5.85), Inches(5.9),
             fill_color=RGBColor(0x0D, 0x1F, 0x1D), line_color=TEAL, line_width_pt=2)

    add_textbox(slide, "✅  GIẢI PHÁP",
                left=Inches(7.28), top=Inches(1.25),
                width=Inches(5.5), height=Inches(0.42),
                font_name="Calibri", font_size=16, bold=True, color=TEAL)

    solutions = [
        ("✅", "Regex Master Normalization",
         "Áp dụng 50+ quy tắc regex chuẩn hóa tên ga: viết thường, loại dấu, alias mapping"),
        ("✅", "Đọc file theo Chunk",
         "pandas read_csv với chunksize=50,000 → xử lý 500K+ dòng mà không tràn RAM"),
        ("✅", "Spatial Join Nearest Neighbor",
         "GeoPandas sjoin_nearest() gán Borough dựa trên polygon địa giới hành chính"),
        ("✅", "Logic Fuzzy + Priority Score",
         "Ưu tiên match chính xác > partial > fuzzy với ngưỡng similarity ≥ 80%"),
        ("✅", "Validation & Logging",
         "Ghi log toàn bộ ga không khớp vào unmatched_stations_log.pdf để kiểm tra"),
    ]

    y = Inches(1.75)
    for icon, title, desc in solutions:
        add_textbox(slide, f"{icon}  {title}",
                    left=Inches(7.28), top=y,
                    width=Inches(5.5), height=Inches(0.32),
                    font_name="Calibri", font_size=13, bold=True,
                    color=RGBColor(0x6E, 0xE7, 0xB7))
        add_textbox(slide, f"     {desc}",
                    left=Inches(7.28), top=y + Inches(0.3),
                    width=Inches(5.5), height=Inches(0.45),
                    font_name="Calibri", font_size=11,
                    color=RGBColor(0x94, 0xA3, 0xB8))
        y += Inches(0.93)


# ───────────────────────────────────────────────────────────────────────────

def slide06_architecture(prs):
    """Slide 6 – Kiến trúc Pipeline ETL"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, LIGHT_GREY)
    header_band(slide, "KIẾN TRÚC PIPELINE ETL END-TO-END",
                "Luồng xử lý dữ liệu hoàn chỉnh từ nguồn thô đến sản phẩm cuối")
    footer_bar(slide)
    slide_number_label(slide, 6)

    stages = [
        {
            "name": "EXTRACT",
            "icon": "📥",
            "color": RGBColor(0x06, 0x95, 0x88),
            "items": ["stations.kml", "TfL_stations.csv", "Stops.csv", "Parser KML", "Pandas read_csv"],
        },
        {
            "name": "TRANSFORM",
            "icon": "⚙️",
            "color": RGBColor(0x29, 0x78, 0xC8),
            "items": ["Regex Normalize", "Borough Spatial Join", "Feature Engineering", "COVID Impact Calc", "Trend Analysis"],
        },
        {
            "name": "LOAD",
            "icon": "💾",
            "color": RGBColor(0x76, 0x44, 0xC2),
            "items": ["SQLite Database", "london_tfl.db", "298 bản ghi sạch", "london_tfl_cleaned.csv", "Excel Export"],
        },
        {
            "name": "SERVE",
            "icon": "🚀",
            "color": RGBColor(0xF0, 0x7B, 0x1F),
            "items": ["KMeans K=6", "Folium Map", "Biểu đồ xu hướng", "Báo cáo PDF", "Flask API"],
        },
    ]

    box_w = Inches(2.9)
    box_h = Inches(5.6)
    gap = Inches(0.25)
    start_x = Inches(0.35)
    box_y = Inches(1.2)

    for i, stage in enumerate(stages):
        cx = start_x + i * (box_w + gap)

        # Main card
        add_rect(slide, cx, box_y, box_w, box_h,
                 fill_color=WHITE, line_color=stage["color"], line_width_pt=2)

        # Header color band
        add_rect(slide, cx, box_y, box_w, Inches(0.85),
                 fill_color=stage["color"])

        # Icon + title in header
        add_textbox(slide, stage["icon"],
                    left=cx + Inches(0.12), top=box_y + Inches(0.06),
                    width=Inches(0.55), height=Inches(0.5),
                    font_name="Segoe UI Emoji", font_size=22, color=WHITE)
        add_textbox(slide, stage["name"],
                    left=cx + Inches(0.65), top=box_y + Inches(0.16),
                    width=box_w - Inches(0.75), height=Inches(0.5),
                    font_name="Calibri", font_size=18, bold=True,
                    color=WHITE, align=PP_ALIGN.LEFT)

        # Step number
        add_textbox(slide, str(i + 1),
                    left=cx + box_w - Inches(0.45), top=box_y + Inches(0.06),
                    width=Inches(0.38), height=Inches(0.38),
                    font_name="Calibri", font_size=22, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)

        # Items
        txBox = slide.shapes.add_textbox(
            cx + Inches(0.18), box_y + Inches(0.95),
            box_w - Inches(0.3), box_h - Inches(1.1))
        tf = txBox.text_frame; tf.word_wrap = True
        first = True
        for item in stage["items"]:
            if first:
                p = tf.paragraphs[0]; first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(8)
            run = p.add_run()
            run.text = f"▶  {item}"
            run.font.name = "Calibri"
            run.font.size = Pt(13)
            run.font.color.rgb = DARK_NAVY

        # Arrow (between stages)
        if i < len(stages) - 1:
            ax = cx + box_w + Inches(0.04)
            add_textbox(slide, "▶",
                        left=ax, top=box_y + box_h / 2 - Inches(0.22),
                        width=Inches(0.2), height=Inches(0.44),
                        font_name="Calibri", font_size=20, bold=True,
                        color=stage["color"], align=PP_ALIGN.CENTER)

    # Bottom tech stack label
    add_textbox(slide, "🔧  Tech Stack:  Python 3.x  |  Pandas  |  GeoPandas  |  Scikit-learn  |  Folium  |  SQLite  |  Matplotlib",
                left=Inches(0.35), top=Inches(7.0),
                width=Inches(12.5), height=Inches(0.25),
                font_name="Calibri", font_size=10.5,
                color=SLATE, align=PP_ALIGN.CENTER)


# ───────────────────────────────────────────────────────────────────────────

def slide07_innovations(prs):
    """Slide 7 – Đổi mới kỹ thuật"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_NAVY)
    header_band(slide, "ĐỔI MỚI KỸ THUẬT",
                "Các giải pháp sáng tạo được áp dụng trong dự án")
    footer_bar(slide)
    slide_number_label(slide, 7)

    innovations = [
        {
            "icon": "🔤",
            "title": "Chuẩn hóa Tên Ga bằng Regex",
            "color": TEAL,
            "points": [
                "50+ quy tắc regex xử lý viết tắt, dấu câu, dấu apostrophe",
                "Alias mapping: 'Heathrow T123' → 'Heathrow Terminals'",
                "Unicode normalization & lowercase pipeline",
                "Tỷ lệ match: >95% tên ga được ghép thành công",
            ],
        },
        {
            "icon": "🗺️",
            "title": "Gán Borough bằng Spatial Join",
            "color": ACCENT_BLU,
            "points": [
                "GeoPandas sjoin_nearest() – Nearest Neighbor Algorithm",
                "Ranh giới hành chính 32 Borough + City of London",
                "Bổ sung Borough cho 100% bản ghi thiếu",
                "Tolerance: 500m – bù đắp sai số tọa độ",
            ],
        },
        {
            "icon": "📊",
            "title": "Feature Engineering Nâng cao",
            "color": ACCENT_YLW,
            "points": [
                "covid_impact_pct: Mức sụt giảm 2019→2020 (%)",
                "recovery_pct: Mức phục hồi 2020→2021 (%)",
                "avg_passengers: Trung bình 5 năm cho clustering",
                "trend_label: Phân loại xu hướng (Giảm mạnh/nhẹ/Tăng/Ổn định)",
            ],
        },
    ]

    card_w = Inches(12.55)
    card_h = Inches(1.82)
    start_x = Inches(0.35)
    start_y = Inches(1.18)

    for i, innov in enumerate(innovations):
        cy = start_y + i * (card_h + Inches(0.18))
        add_rect(slide, start_x, cy, card_w, card_h,
                 fill_color=MID_NAVY, line_color=innov["color"], line_width_pt=2)

        # Left color strip
        add_rect(slide, start_x, cy, Inches(0.09), card_h,
                 fill_color=innov["color"])

        # Icon
        add_textbox(slide, innov["icon"],
                    left=start_x + Inches(0.2), top=cy + Inches(0.22),
                    width=Inches(0.6), height=Inches(0.55),
                    font_name="Segoe UI Emoji", font_size=26, color=innov["color"])

        # Title
        add_textbox(slide, innov["title"],
                    left=start_x + Inches(0.9), top=cy + Inches(0.2),
                    width=Inches(3.5), height=Inches(0.45),
                    font_name="Calibri", font_size=15, bold=True, color=innov["color"])

        # Points (2 columns of 2)
        points = innov["points"]
        col_w2 = Inches(4.0)
        for j, pt in enumerate(points):
            col = j % 2
            row = j // 2
            px = start_x + Inches(0.9) + col * col_w2
            py = cy + Inches(0.72) + row * Inches(0.44)
            add_textbox(slide, f"▸  {pt}",
                        left=px, top=py,
                        width=col_w2 - Inches(0.1), height=Inches(0.4),
                        font_name="Calibri", font_size=12,
                        color=RGBColor(0xCB, 0xD5, 0xE1))


# ───────────────────────────────────────────────────────────────────────────

def slide08_kmeans(prs):
    """Slide 8 – Phân cụm KMeans"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, LIGHT_GREY)
    header_band(slide, "PHÂN CỤM KMEANS  (K = 6)",
                "Phân loại 298 nhà ga theo hành vi lưu lượng hành khách")
    footer_bar(slide)
    slide_number_label(slide, 8)

    # Left: cluster table
    clusters = [
        ("Siêu trung tâm", "8 ga",  "37.32M", "-73.70%", RGBColor(0x06, 0xB6, 0xD4)),
        ("Ga lớn",         "42 ga", "11.27M", "-63.47%", RGBColor(0x38, 0x78, 0xE8)),
        ("Ga trung bình",  "78 ga", "3.14M",  "-54.91%", RGBColor(0xF9, 0x73, 0x16)),
        ("Ga nhỏ",         "80 ga", "2.83M",  "-57.42%", RGBColor(0xEC, 0x48, 0x99)),
        ("Ga ít khách",    "37 ga", "2.22M",  "-49.69%", RGBColor(0x22, 0xC5, 0x5E)),
        ("Ga rất ít khách","53 ga", "1.63M",  "-41.42%", RGBColor(0xA8, 0x55, 0xF7)),
    ]

    table_x = Inches(0.32)
    table_y = Inches(1.18)
    table_w = Inches(5.65)
    col_widths = [Inches(1.72), Inches(0.85), Inches(1.25), Inches(1.25), Inches(0.58)]
    headers = ["Cụm", "Số ga", "Khách TB 2021", "COVID Impact", ""]
    row_h = Inches(0.52)

    # Header row
    x = table_x
    for h, cw in zip(headers, col_widths):
        add_rect(slide, x, table_y, cw, row_h, fill_color=DARK_NAVY)
        add_textbox(slide, h,
                    left=x + Inches(0.06), top=table_y + Inches(0.1),
                    width=cw - Inches(0.1), height=row_h - Inches(0.12),
                    font_name="Calibri", font_size=11, bold=True,
                    color=TEAL, align=PP_ALIGN.CENTER)
        x += cw

    # Data rows
    for ri, (name, count, avg, covid, color) in enumerate(clusters):
        ry = table_y + (ri + 1) * row_h
        bg = WHITE if ri % 2 == 0 else LIGHT_GREY
        x = table_x
        row_data = [name, count, avg, covid, ""]
        for ci, (val, cw) in enumerate(zip(row_data, col_widths)):
            add_rect(slide, x, ry, cw, row_h, fill_color=bg,
                     line_color=RGBColor(0xE2, 0xE8, 0xF0), line_width_pt=0.5)
            if ci == 0:
                # Color dot + name
                dot = slide.shapes.add_shape(9, x + Inches(0.08), ry + Inches(0.17),
                                              Inches(0.16), Inches(0.18))
                dot.fill.solid(); dot.fill.fore_color.rgb = color
                dot.line.fill.background()
                add_textbox(slide, name,
                            left=x + Inches(0.3), top=ry + Inches(0.1),
                            width=cw - Inches(0.35), height=row_h - Inches(0.12),
                            font_name="Calibri", font_size=12, bold=True,
                            color=DARK_NAVY)
            elif ci == 2:  # Avg passengers
                add_textbox(slide, val,
                            left=x + Inches(0.06), top=ry + Inches(0.1),
                            width=cw - Inches(0.1), height=row_h - Inches(0.12),
                            font_name="Calibri", font_size=13, bold=True,
                            color=RGBColor(0x0F, 0x17, 0x2A), align=PP_ALIGN.CENTER)
            elif ci == 3:  # COVID
                add_textbox(slide, val,
                            left=x + Inches(0.06), top=ry + Inches(0.1),
                            width=cw - Inches(0.1), height=row_h - Inches(0.12),
                            font_name="Calibri", font_size=12, bold=True,
                            color=RGBColor(0xDC, 0x26, 0x26), align=PP_ALIGN.CENTER)
            else:
                add_textbox(slide, val,
                            left=x + Inches(0.06), top=ry + Inches(0.1),
                            width=cw - Inches(0.1), height=row_h - Inches(0.12),
                            font_name="Calibri", font_size=12,
                            color=DARK_NAVY, align=PP_ALIGN.CENTER)
            x += cw

    # Features used label
    add_textbox(slide, "📐  Features: avg_passengers, covid_impact_pct, recovery_pct  |  Algorithm: KMeans (sklearn)  |  Elbow Method → K=6",
                left=Inches(0.32), top=table_y + (len(clusters) + 1) * row_h + Inches(0.1),
                width=table_w, height=Inches(0.32),
                font_name="Calibri", font_size=9.5, italic=True,
                color=SLATE)

    # Right: chart
    add_rect(slide, Inches(6.2), Inches(1.18), Inches(6.78), Inches(5.9),
             fill_color=WHITE, line_color=TEAL, line_width_pt=1.5)
    add_textbox(slide, "Biểu đồ so sánh theo Cụm",
                left=Inches(6.38), top=Inches(1.22),
                width=Inches(6.4), height=Inches(0.3),
                font_name="Calibri", font_size=11, italic=True,
                color=SLATE)
    add_image(slide, IMG_CLUSTER,
              Inches(6.22), Inches(1.52), Inches(6.74), Inches(5.52))


# ───────────────────────────────────────────────────────────────────────────

def slide09_covid(prs):
    """Slide 9 – Phân tích tác động COVID-19"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_NAVY)
    header_band(slide, "PHÂN TÍCH TÁC ĐỘNG COVID-19",
                "Biến động lưu lượng hành khách 2017 – 2021")
    footer_bar(slide)
    slide_number_label(slide, 9)

    # Chart (left 2/3)
    add_rect(slide, Inches(0.32), Inches(1.18), Inches(8.2), Inches(5.88),
             fill_color=WHITE, line_color=TEAL, line_width_pt=1.5)
    add_image(slide, IMG_TRENDS,
              Inches(0.34), Inches(1.2), Inches(8.16), Inches(5.84))

    # Right stats panel
    rx = Inches(8.75)
    add_textbox(slide, "📊  SỐ LIỆU CHÍNH",
                left=rx, top=Inches(1.25),
                width=Inches(4.2), height=Inches(0.4),
                font_name="Calibri", font_size=14, bold=True, color=TEAL)

    kpis = [
        ("3.28 Tỷ", "Tổng hành khách 2019\n(Cao nhất 5 năm)"),
        ("-63.21%", "Sụt giảm 2019→2020\n(Tác động COVID)"),
        ("1.21 Tỷ", "Tổng hành khách 2020\n(Thấp nhất lịch sử)"),
        ("+16.82%", "Phục hồi 2020→2021"),
    ]
    colors_kpi = [WHITE, RGBColor(0xFF, 0x6B, 0x6B), RGBColor(0xFF, 0xA5, 0xA5), RGBColor(0x6E, 0xE7, 0xB7)]
    ky = Inches(1.75)
    for (val, lbl), col in zip(kpis, colors_kpi):
        add_rect(slide, rx, ky, Inches(4.22), Inches(1.12),
                 fill_color=MID_NAVY, line_color=col, line_width_pt=1.5)
        add_textbox(slide, val,
                    left=rx + Inches(0.15), top=ky + Inches(0.08),
                    width=Inches(3.9), height=Inches(0.55),
                    font_name="Calibri", font_size=26, bold=True, color=col)
        add_textbox(slide, lbl,
                    left=rx + Inches(0.15), top=ky + Inches(0.6),
                    width=Inches(3.9), height=Inches(0.45),
                    font_name="Calibri", font_size=11,
                    color=RGBColor(0x94, 0xA3, 0xB8))
        ky += Inches(1.2)

    # Key insight
    add_rect(slide, rx, ky + Inches(0.1), Inches(4.22), Inches(0.88),
             fill_color=RGBColor(0x0D, 0x1F, 0x1D), line_color=TEAL, line_width_pt=1.2)
    add_textbox(slide, "💡  287/298 ga có xu hướng\ngiảm mạnh – chỉ 1 ga tăng trưởng",
                left=rx + Inches(0.15), top=ky + Inches(0.15),
                width=Inches(3.9), height=Inches(0.7),
                font_name="Calibri", font_size=12, italic=True,
                color=TEAL)


# ───────────────────────────────────────────────────────────────────────────

def slide10_results(prs):
    """Slide 10 – Kết quả nổi bật"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, LIGHT_GREY)
    header_band(slide, "KẾT QUẢ NỔI BẬT",
                "Tổng hợp các phát hiện quan trọng từ phân tích 298 nhà ga")
    footer_bar(slide)
    slide_number_label(slide, 10)

    # Top KPI row
    kpis = [
        ("298", "Nhà ga phân tích", "🚉"),
        ("6", "Cụm hành vi", "🤖"),
        ("-63.21%", "Sụt giảm COVID", "🦠"),
        ("+16.82%", "Phục hồi 2021", "📈"),
    ]
    kpi_w = Inches(3.0)
    kpi_h = Inches(1.15)
    kpi_y = Inches(1.18)
    gap_k = Inches(0.2)
    start_kpi_x = Inches(0.35)
    for i, (val, lbl, icon) in enumerate(kpis):
        kx = start_kpi_x + i * (kpi_w + gap_k)
        col = RGBColor(0xFF, 0x6B, 0x6B) if "-" in val else TEAL
        add_rect(slide, kx, kpi_y, kpi_w, kpi_h,
                 fill_color=DARK_NAVY, line_color=col, line_width_pt=2)
        add_textbox(slide, icon,
                    left=kx + Inches(0.12), top=kpi_y + Inches(0.1),
                    width=Inches(0.45), height=Inches(0.45),
                    font_name="Segoe UI Emoji", font_size=20, color=col)
        add_textbox(slide, val,
                    left=kx + Inches(0.55), top=kpi_y + Inches(0.08),
                    width=kpi_w - Inches(0.65), height=Inches(0.55),
                    font_name="Calibri", font_size=22, bold=True, color=col)
        add_textbox(slide, lbl,
                    left=kx + Inches(0.12), top=kpi_y + Inches(0.62),
                    width=kpi_w - Inches(0.2), height=Inches(0.38),
                    font_name="Calibri", font_size=11,
                    color=RGBColor(0x94, 0xA3, 0xB8))

    # Bottom section: Top 5 + Cluster distribution
    # Left: Top 5 stations
    add_rect(slide, Inches(0.35), Inches(2.45), Inches(6.2), Inches(4.6),
             fill_color=WHITE, line_color=TEAL, line_width_pt=1.5)
    add_textbox(slide, "🏆  TOP 5 GA ĐÔNG KHÁCH NHẤT (2021)",
                left=Inches(0.53), top=Inches(2.55),
                width=Inches(5.8), height=Inches(0.38),
                font_name="Calibri", font_size=13, bold=True, color=DARK_NAVY)

    top5 = [
        ("1", "King's Cross St. Pancras", "34.0M", 1.0),
        ("2", "Waterloo",                 "29.7M", 0.87),
        ("3", "Victoria",                 "28.3M", 0.83),
        ("4", "London Bridge",            "27.2M", 0.80),
        ("5", "Bank / Monument",          "26.1M", 0.77),
    ]
    bar_max_w = Inches(3.4)
    ty = Inches(3.05)
    for rank, name, count, ratio in top5:
        add_textbox(slide, rank,
                    left=Inches(0.5), top=ty,
                    width=Inches(0.28), height=Inches(0.38),
                    font_name="Calibri", font_size=13, bold=True, color=TEAL)
        add_textbox(slide, name,
                    left=Inches(0.82), top=ty,
                    width=Inches(2.2), height=Inches(0.38),
                    font_name="Calibri", font_size=12, color=DARK_NAVY)
        # Bar
        add_rect(slide, Inches(3.1), ty + Inches(0.07),
                 bar_max_w * ratio, Inches(0.24), fill_color=TEAL)
        add_textbox(slide, count,
                    left=Inches(3.1) + bar_max_w * ratio + Inches(0.05), top=ty,
                    width=Inches(0.7), height=Inches(0.38),
                    font_name="Calibri", font_size=11, bold=True,
                    color=DARK_NAVY)
        ty += Inches(0.72)

    # Right: cluster distribution
    add_rect(slide, Inches(6.75), Inches(2.45), Inches(6.18), Inches(4.6),
             fill_color=WHITE, line_color=ACCENT_BLU, line_width_pt=1.5)
    add_textbox(slide, "📊  PHÂN BỐ THEO CỤM",
                left=Inches(6.93), top=Inches(2.55),
                width=Inches(5.8), height=Inches(0.38),
                font_name="Calibri", font_size=13, bold=True, color=DARK_NAVY)

    dist = [
        ("Siêu trung tâm",  8,  298, RGBColor(0x06, 0xB6, 0xD4)),
        ("Ga lớn",         42,  298, RGBColor(0x38, 0x78, 0xE8)),
        ("Ga trung bình",  78,  298, RGBColor(0xF9, 0x73, 0x16)),
        ("Ga nhỏ",         80,  298, RGBColor(0xEC, 0x48, 0x99)),
        ("Ga ít khách",    37,  298, RGBColor(0x22, 0xC5, 0x5E)),
        ("Ga rất ít khách",53,  298, RGBColor(0xA8, 0x55, 0xF7)),
    ]
    bar_max = Inches(3.5)
    dy = Inches(3.05)
    for name, count, total, color in dist:
        ratio = count / total
        add_textbox(slide, name,
                    left=Inches(6.93), top=dy,
                    width=Inches(1.85), height=Inches(0.35),
                    font_name="Calibri", font_size=11, color=DARK_NAVY)
        add_rect(slide, Inches(8.82), dy + Inches(0.05),
                 bar_max * ratio, Inches(0.22), fill_color=color)
        add_textbox(slide, str(count),
                    left=Inches(8.82) + bar_max * ratio + Inches(0.06), top=dy,
                    width=Inches(0.4), height=Inches(0.35),
                    font_name="Calibri", font_size=11, bold=True, color=DARK_NAVY)
        dy += Inches(0.62)


# ───────────────────────────────────────────────────────────────────────────

def slide11_products(prs):
    """Slide 11 – Sản phẩm & Trực quan hóa"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_NAVY)
    header_band(slide, "SẢN PHẨM & TRỰC QUAN HÓA",
                "Các đầu ra được tạo ra từ Pipeline ETL")
    footer_bar(slide)
    slide_number_label(slide, 11)

    products = [
        {
            "icon": "🗺️",
            "title": "Bản đồ Tương tác Folium",
            "file": "london_tfl_map.html",
            "color": TEAL,
            "desc": [
                "298 marker nhà ga với popup thông tin đầy đủ",
                "Color-coded theo Cluster (6 màu phân biệt)",
                "Hover hiển thị: tên ga, Borough, lượng khách 2021",
                "Zoom tương tác – chạy trực tiếp trên trình duyệt",
            ],
        },
        {
            "icon": "📄",
            "title": "Báo cáo PDF Tự động",
            "file": "TfL_Project_Report.pdf",
            "color": ACCENT_BLU,
            "desc": [
                "Sinh tự động bằng ReportLab Python",
                "Bao gồm: tóm tắt, bảng cluster, biểu đồ",
                "Thống kê Borough & phân tích xu hướng",
                "Export ngay sau khi pipeline hoàn thành",
            ],
        },
        {
            "icon": "🗄️",
            "title": "SQLite Database",
            "file": "london_tfl.db",
            "color": ACCENT_YLW,
            "desc": [
                "Lưu trữ 298 bản ghi sạch hoàn chỉnh",
                "Schema: station_id, name, borough, coordinates",
                "Passenger data 2017-2021 + cluster labels",
                "Hỗ trợ truy vấn SQL trực tiếp",
            ],
        },
        {
            "icon": "📊",
            "title": "Excel / CSV Export",
            "file": "london_tfl_results.xlsx",
            "color": RGBColor(0x22, 0xC5, 0x5E),
            "desc": [
                "london_tfl_cleaned.csv – dữ liệu sạch",
                "london_tfl_results.xlsx – kết quả đầy đủ",
                "Sẵn sàng import vào Power BI, Tableau",
                "Unmatched log PDF cho kiểm tra chất lượng",
            ],
        },
    ]

    card_w = Inches(6.05)
    card_h = Inches(2.55)
    start_x = Inches(0.32)
    start_y = Inches(1.18)

    for i, prod in enumerate(products):
        col = i % 2
        row = i // 2
        cx = start_x + col * (card_w + Inches(0.25))
        cy = start_y + row * (card_h + Inches(0.22))

        add_rect(slide, cx, cy, card_w, card_h,
                 fill_color=MID_NAVY, line_color=prod["color"], line_width_pt=1.8)

        # Left color strip
        add_rect(slide, cx, cy, Inches(0.08), card_h, fill_color=prod["color"])

        # Icon
        add_textbox(slide, prod["icon"],
                    left=cx + Inches(0.2), top=cy + Inches(0.18),
                    width=Inches(0.55), height=Inches(0.5),
                    font_name="Segoe UI Emoji", font_size=24, color=prod["color"])

        # Title + file
        add_textbox(slide, prod["title"],
                    left=cx + Inches(0.85), top=cy + Inches(0.12),
                    width=card_w - Inches(0.95), height=Inches(0.38),
                    font_name="Calibri", font_size=14, bold=True, color=prod["color"])
        add_textbox(slide, f"📁  {prod['file']}",
                    left=cx + Inches(0.85), top=cy + Inches(0.5),
                    width=card_w - Inches(0.95), height=Inches(0.28),
                    font_name="Calibri", font_size=10, italic=True,
                    color=RGBColor(0x64, 0x74, 0x8B))

        # Bullet points
        txBox = slide.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.9),
                                          card_w - Inches(0.3), card_h - Inches(1.0))
        tf = txBox.text_frame; tf.word_wrap = True
        first = True
        for pt in prod["desc"]:
            if first:
                p = tf.paragraphs[0]; first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = f"▸  {pt}"
            run.font.name = "Calibri"
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)


# ───────────────────────────────────────────────────────────────────────────

def slide12_conclusion(prs):
    """Slide 12 – Kết luận & Hướng phát triển"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, LIGHT_GREY)
    header_band(slide, "KẾT LUẬN & HƯỚNG PHÁT TRIỂN",
                "Thành tựu đạt được và lộ trình cải tiến trong tương lai")
    footer_bar(slide)
    slide_number_label(slide, 12)

    # Left: conclusions
    add_rect(slide, Inches(0.32), Inches(1.18), Inches(5.9), Inches(5.88),
             fill_color=WHITE, line_color=TEAL, line_width_pt=2)

    add_textbox(slide, "✅  KẾT QUẢ ĐẠT ĐƯỢC",
                left=Inches(0.5), top=Inches(1.3),
                width=Inches(5.5), height=Inches(0.4),
                font_name="Calibri", font_size=15, bold=True, color=TEAL)

    conclusions = [
        ("🔧", "Pipeline ETL hoàn chỉnh",
         "Tự động hóa toàn bộ Extract→Transform→Load→Serve"),
        ("🗄️", "Dữ liệu sạch & chuẩn hóa",
         "298 ga với đầy đủ tọa độ, Borough, lưu lượng 5 năm"),
        ("🤖", "Phân cụm KMeans K=6",
         "Phân loại hành vi 298 ga, xác định cluster Siêu trung tâm"),
        ("🦠", "Phân tích COVID-19",
         "Định lượng mức sụt giảm -63.21% và phục hồi +16.82%"),
        ("🗺️", "Trực quan hóa tương tác",
         "Bản đồ Folium, biểu đồ và báo cáo PDF tự động"),
    ]

    cy = Inches(1.82)
    for icon, title, desc in conclusions:
        add_rect(slide, Inches(0.5), cy, Inches(5.5), Inches(0.82),
                 fill_color=RGBColor(0xF8, 0xFF, 0xFE),
                 line_color=RGBColor(0xE0, 0xF5, 0xF3), line_width_pt=1)
        add_textbox(slide, icon,
                    left=Inches(0.55), top=cy + Inches(0.12),
                    width=Inches(0.38), height=Inches(0.38),
                    font_name="Segoe UI Emoji", font_size=18, color=TEAL)
        add_textbox(slide, title,
                    left=Inches(0.97), top=cy + Inches(0.08),
                    width=Inches(4.8), height=Inches(0.32),
                    font_name="Calibri", font_size=13, bold=True, color=DARK_NAVY)
        add_textbox(slide, desc,
                    left=Inches(0.97), top=cy + Inches(0.4),
                    width=Inches(4.8), height=Inches(0.35),
                    font_name="Calibri", font_size=11, color=SLATE)
        cy += Inches(0.95)

    # Right: future development
    add_rect(slide, Inches(6.42), Inches(1.18), Inches(6.5), Inches(5.88),
             fill_color=DARK_NAVY, line_color=ACCENT_BLU, line_width_pt=2)

    add_textbox(slide, "🚀  HƯỚNG PHÁT TRIỂN",
                left=Inches(6.6), top=Inches(1.3),
                width=Inches(6.1), height=Inches(0.4),
                font_name="Calibri", font_size=15, bold=True, color=ACCENT_BLU)

    future = [
        ("🔤", "Fuzzy Matching nâng cao",
         "Áp dụng RapidFuzz / FuzzyWuzzy để cải thiện match rate tên ga"),
        ("⚡", "Apache Kafka Streaming",
         "Xử lý dữ liệu thời gian thực thay vì batch processing"),
        ("🌊", "Apache Airflow Orchestration",
         "Lên lịch tự động, monitoring pipeline với DAG"),
        ("☁️", "Cloud Migration",
         "Chuyển lên AWS S3 + Redshift hoặc GCP BigQuery"),
        ("🤖", "Machine Learning nâng cao",
         "Dự báo lưu lượng, phát hiện bất thường (LSTM, Prophet)"),
        ("📱", "Dashboard Real-time",
         "Power BI / Grafana cho giám sát vận hành liên tục"),
    ]

    fy = Inches(1.82)
    for icon, title, desc in future:
        add_textbox(slide, f"{icon}  {title}",
                    left=Inches(6.6), top=fy,
                    width=Inches(6.1), height=Inches(0.32),
                    font_name="Calibri", font_size=13, bold=True, color=TEAL)
        add_textbox(slide, desc,
                    left=Inches(6.6), top=fy + Inches(0.3),
                    width=Inches(6.1), height=Inches(0.38),
                    font_name="Calibri", font_size=11,
                    color=RGBColor(0x94, 0xA3, 0xB8))
        fy += Inches(0.9)


# ───────────────────────────────────────────────────────────────────────────

def slide13_thankyou(prs):
    """Slide 13 – Cảm ơn & Q&A"""
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_NAVY)

    # Top teal stripe
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.09), fill_color=TEAL)

    # Bottom teal stripe
    add_rect(slide, 0, SLIDE_H - Inches(0.09), SLIDE_W, Inches(0.09), fill_color=TEAL)

    # Background decorative circles
    for radius, alpha_color, lx, ly in [
        (Inches(4.5), RGBColor(0x00, 0x40, 0x38), Inches(-1.5), Inches(-1.5)),
        (Inches(3.0), RGBColor(0x0A, 0x29, 0x45), Inches(11), Inches(5)),
        (Inches(2.0), RGBColor(0x1E, 0x29, 0x3B), Inches(1), Inches(5.5)),
    ]:
        circ = slide.shapes.add_shape(9, lx, ly, radius, radius)
        circ.fill.solid(); circ.fill.fore_color.rgb = alpha_color
        circ.line.fill.background()

    # Main content
    add_textbox(slide, "CẢM ƠN QUÝ THẦY CÔ & CÁC BẠN ĐÃ LẮNG NGHE!",
                left=Inches(1.0), top=Inches(1.4),
                width=Inches(11.33), height=Inches(0.85),
                font_name="Calibri", font_size=30, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    # Teal accent
    add_rect(slide, Inches(5.0), Inches(2.35), Inches(3.33), Inches(0.07), fill_color=TEAL)

    # Q&A
    add_textbox(slide, "❓  Q & A",
                left=Inches(1.0), top=Inches(2.55),
                width=Inches(11.33), height=Inches(0.65),
                font_name="Calibri", font_size=38, bold=True,
                color=TEAL, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Xin mời thầy cô và các bạn đặt câu hỏi",
                left=Inches(1.0), top=Inches(3.25),
                width=Inches(11.33), height=Inches(0.38),
                font_name="Calibri", font_size=16, italic=True,
                color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, Inches(1.5), Inches(3.75), Inches(10.33), Inches(0.04),
             fill_color=RGBColor(0x1E, 0x29, 0x3B))

    # Group info block
    add_rect(slide, Inches(0.8), Inches(3.95), Inches(11.73), Inches(2.85),
             fill_color=MID_NAVY, line_color=TEAL, line_width_pt=1.5)

    add_textbox(slide, "THÔNG TIN NHÓM",
                left=Inches(1.0), top=Inches(4.08),
                width=Inches(11.33), height=Inches(0.4),
                font_name="Calibri", font_size=14, bold=True,
                color=TEAL, align=PP_ALIGN.CENTER)

    # Group details
    group_info = [
        ("👥", "Nhóm 14  |  Lớp DHKHDL20A  |  Khoa Khoa học Dữ liệu"),
        ("👨‍🏫", "Giảng viên hướng dẫn: TS. Lê Trọng Ngọc"),
        ("📚", "Môn học: Data Engineering  |  Học kỳ 2, Năm học 2024–2025"),
        ("🔗", "Dự án: Pipeline ETL End-to-End – Hệ thống TfL London (298 nhà ga)"),
    ]
    gy = Inches(4.58)
    for icon, text in group_info:
        add_textbox(slide, f"{icon}   {text}",
                    left=Inches(1.5), top=gy,
                    width=Inches(10.33), height=Inches(0.34),
                    font_name="Calibri", font_size=13,
                    color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
        gy += Inches(0.42)

    slide_number_label(slide, 13)


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    print("🚀 Đang tạo file PowerPoint...")
    prs = new_prs()

    print("  ✅ Slide  1: Trang bìa")
    slide01_cover(prs)

    print("  ✅ Slide  2: Lời mở đầu & Tóm tắt")
    slide02_intro(prs)

    print("  ✅ Slide  3: Lý do chọn đề tài")
    slide03_reasons(prs)

    print("  ✅ Slide  4: Bộ dữ liệu sử dụng")
    slide04_datasets(prs)

    print("  ✅ Slide  5: Thách thức & Giải pháp")
    slide05_challenges(prs)

    print("  ✅ Slide  6: Kiến trúc Pipeline ETL")
    slide06_architecture(prs)

    print("  ✅ Slide  7: Đổi mới kỹ thuật")
    slide07_innovations(prs)

    print("  ✅ Slide  8: Phân cụm KMeans")
    slide08_kmeans(prs)

    print("  ✅ Slide  9: Phân tích COVID-19")
    slide09_covid(prs)

    print("  ✅ Slide 10: Kết quả nổi bật")
    slide10_results(prs)

    print("  ✅ Slide 11: Sản phẩm & Trực quan hóa")
    slide11_products(prs)

    print("  ✅ Slide 12: Kết luận & Hướng phát triển")
    slide12_conclusion(prs)

    print("  ✅ Slide 13: Cảm ơn & Q&A")
    slide13_thankyou(prs)

    prs.save(OUT_PPTX)
    print(f"\n🎉 Hoàn thành! File đã được lưu tại:\n   {OUT_PPTX}")
    size_mb = os.path.getsize(OUT_PPTX) / (1024 * 1024)
    print(f"   📦 Kích thước: {size_mb:.2f} MB")


if __name__ == "__main__":
    main()
