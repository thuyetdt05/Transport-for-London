# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
prs = Presentation('Nhom14_TfL_ETL_Pipeline.pptx')

print('=== SLIDE 7 — KY THUAT (sau sua) ===')
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            keywords = ['Spatial', 'Nearest', 'Pandas', 'Tolerance', 'StopType', 'geopandas', 'sjoin']
            if any(k.lower() in t.lower() for k in keywords):
                print(f'  [{shape.name}] {t}')

print()
print('=== SLIDE 8 — KMEANS FEATURES (sau sua) ===')
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if shape.name == 'TextBox 84' and shape.has_text_frame:
        print(f'  [{shape.name}] {shape.text_frame.paragraphs[0].text}')

print()
print('=== SLIDE 10 — TOP 5 GA (sau sua) ===')
slide10 = prs.slides[9]
top5_pairs = [
    ('TextBox 27', 'TextBox 29'),
    ('TextBox 31', 'TextBox 33'),
    ('TextBox 35', 'TextBox 37'),
    ('TextBox 39', 'TextBox 41'),
    ('TextBox 43', 'TextBox 45'),
]
shape_map = {}
for shape in slide10.shapes:
    if shape.has_text_frame:
        shape_map[shape.name] = shape.text_frame.paragraphs[0].text

for i, (n, v) in enumerate(top5_pairs, 1):
    ga = shape_map.get(n, '?')
    hk = shape_map.get(v, '?')
    print(f'  #{i}: {ga}  --  {hk}')
